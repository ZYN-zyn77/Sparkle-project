import os
import re
from typing import List, Dict, Optional
from fastapi import HTTPException
from loguru import logger
from pydantic import BaseModel

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    Document = None
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    Presentation = None
    HAS_PPTX = False

try:
    import pytesseract
    from PIL import Image, ImageOps, ImageEnhance
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

class ExtractedChunk(BaseModel):
    text: str
    page_num: int
    source: str
    metadata: Dict = {}  # bold, header, color, etc.
    ocr_confidence: Optional[float] = None  # 0.0-1.0, None if not from OCR

class IngestionService:
    """
    Robust Document Ingestion Service for Exam Savior.
    Handles PDF, DOCX, PPTX with advanced cleaning and metadata extraction.
    """

    def process_file(self, file_path: str) -> List[ExtractedChunk]:
        """
        Main entry point. Dispatches to specific handlers based on extension.
        Includes Magic Byte validation.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        self._validate_magic_bytes(file_path)

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        try:
            if ext == ".pdf":
                return self._process_pdf(file_path)
            elif ext == ".docx":
                return self._process_docx(file_path)
            elif ext == ".pptx":
                return self._process_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return []
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            raise e

    def _validate_magic_bytes(self, file_path: str):
        """
        Check if the file header matches its extension to prevent spoofing.
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        with open(file_path, "rb") as f:
            header = f.read(4)

        if ext == ".pdf":
            if header != b"%PDF":
                raise ValueError(f"Invalid PDF header: {header}")
        elif ext in [".docx", ".pptx"]:
            if header != b"PK\x03\x04":
                raise ValueError(f"Invalid ZIP/Office header: {header}")

    def _process_pdf(self, path: str) -> List[ExtractedChunk]:
        if not HAS_PDFPLUMBER:
            raise HTTPException(
                status_code=501,
                detail="PDF processing requires pdfplumber, which is not installed."
            )
        chunks = []
        # Use pdfplumber for better layout analysis
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                ocr_confidence = None

                # --- OCR Fallback Strategy ---
                # If text is empty or suspiciously short (scanned page), try OCR
                if len(text.strip()) < 50:
                    logger.info(f"Page {i+1} has low text content ({len(text.strip())} chars). Attempting OCR...")
                    ocr_text, ocr_confidence = self._attempt_ocr(page)
                    if ocr_text:
                        text = ocr_text
                        logger.info(
                            f"OCR recovered {len(text)} chars from Page {i+1} "
                            f"(confidence: {ocr_confidence:.2f})" if ocr_confidence else
                            f"OCR recovered {len(text)} chars from Page {i+1}"
                        )
                    else:
                        logger.warning(f"OCR failed or produced no text for Page {i+1}")

                if not text:
                    continue

                # Cleaning
                clean_text = self._clean_text(text)

                if len(clean_text) < 20:  # Skip empty/noise pages
                    continue

                chunks.append(ExtractedChunk(
                    text=clean_text,
                    page_num=i + 1,
                    source="pdf",
                    metadata={"raw_len": len(text)},
                    ocr_confidence=ocr_confidence
                ))
        return chunks

    def _attempt_ocr(self, page) -> tuple[str, Optional[float]]:
        """
        Helper to run Tesseract on a pdfplumber page object with image preprocessing.

        Returns:
            tuple: (extracted_text, confidence_score)
                - confidence_score is 0.0-1.0, None if OCR failed
        """
        if not HAS_OCR:
            logger.warning("OCR requested but pytesseract/Pillow not installed.")
            return "", None

        try:
            # Convert page to image (Resolution 300 DPI for better OCR)
            im = page.to_image(resolution=300).original

            # --- Image Preprocessing for Accuracy ---
            # 1. Convert to grayscale
            im = im.convert('L')

            # 2. Auto-contrast (stretch histogram)
            im = ImageOps.autocontrast(im)

            # 3. Simple Binarization (Thresholding)
            threshold = 200
            im = im.point(lambda p: 255 if p > threshold else 0)

            # Run OCR with confidence data
            try:
                config = r'--oem 3 --psm 6'
                # Get detailed data including confidence
                ocr_data = pytesseract.image_to_data(
                    im, lang='chi_sim+eng', config=config, output_type=pytesseract.Output.DICT
                )
            except pytesseract.TesseractError:
                config = r'--oem 3 --psm 6'
                ocr_data = pytesseract.image_to_data(
                    im, lang='eng', config=config, output_type=pytesseract.Output.DICT
                )

            # Extract text and calculate average confidence
            texts = []
            confidences = []

            for i, text in enumerate(ocr_data.get('text', [])):
                text = text.strip()
                conf = ocr_data.get('conf', [])[i]

                if text and conf != -1:  # -1 means no confidence available
                    texts.append(text)
                    confidences.append(conf)

            full_text = ' '.join(texts)

            # Calculate average confidence (0-100 scale -> 0-1 scale)
            avg_confidence = None
            if confidences:
                avg_confidence = sum(confidences) / len(confidences) / 100.0

            return full_text, avg_confidence

        except Exception as e:
            logger.warning(f"OCR Error: {e}")
            return "", None

    def _process_docx(self, path: str) -> List[ExtractedChunk]:
        if not HAS_DOCX:
            raise HTTPException(
                status_code=501,
                detail="DOCX processing requires python-docx, which is not installed."
            )
        doc = Document(path)
        chunks = []
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue

            # Feature Engineering: Extract styles
            style_name = para.style.name.lower()
            is_header = "heading" in style_name
            
            # Check for bold/color runs
            is_bold = any(run.bold for run in para.runs)
            
            metadata = {
                "is_header": is_header,
                "is_bold": is_bold,
                "style": style_name
            }

            clean_text = self._clean_text(text)
            
            chunks.append(ExtractedChunk(
                text=clean_text,
                page_num=i, # Docx doesn't have strict pages, use para index
                source="docx",
                metadata=metadata
            ))
        return chunks

    def _process_pptx(self, path: str) -> List[ExtractedChunk]:
        if not HAS_PPTX:
            raise HTTPException(
                status_code=501,
                detail="PPTX processing requires python-pptx, which is not installed."
            )
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            # CRITICAL: Extract speaker notes (where the real content often lives)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    slide_text.append(f"[NOTES]: {notes}")

            full_text = "\n".join(slide_text)
            clean_text = self._clean_text(full_text)

            if len(clean_text) < 10:
                continue

            chunks.append(ExtractedChunk(
                text=clean_text,
                page_num=i + 1,
                source="pptx",
                metadata={}
            ))
        return chunks

    def _clean_text(self, text: str) -> str:
        """
        Applies cleaning rules defined in Protocol.
        """
        # 1. De-hyphenation (simple heuristic)
        # Fix "exam-\nple" -> "example"
        text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)

        # 2. Remove Page Numbers (standalone numbers on lines)
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)

        # 3. Collapse multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)

        # 4. Remove common watermarks (example)
        text = text.replace("Do Not Distribute", "")
        
        # 5. Fix common encoding artifacts
        text = text.replace("\x00", "") # Null bytes

        return text.strip()

# Singleton instance
ingestion_service = IngestionService()
